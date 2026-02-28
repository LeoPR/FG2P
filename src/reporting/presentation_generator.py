#!/usr/bin/env python
"""
Gerador de apresentação PPTX para FG2P — Unified Mode System
Conteúdo baseado em docs/17_APRESENTACAO.md com tags [modes: full|compact]

O arquivo MD único controla quais slides aparecem em cada versão através de tags.
Sintaxe: [modes: full, compact] — slide aparece em ambas
         [modes: full]           — slide só na versão full
         [modes: compact]        — slide só na versão compact

Uso:
    python src/reporting/presentation_generator.py --mode full        # 29 slides completos
    python src/reporting/presentation_generator.py --mode compact     # 16 slides (10 min)
    python src/reporting/presentation_generator.py --mode full -o results/fg2p.pptx
    python src/reporting/presentation_generator.py                    # default: mode=full
"""

import argparse
import sys
from pathlib import Path

sys.path.append(str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree

try:
    from utils import get_logger, RESULTS_DIR
except ImportError:
    import logging
    def get_logger(name):
        logging.basicConfig(level=logging.INFO)
        return logging.getLogger(name)
    RESULTS_DIR = Path(__file__).parent.parent.parent / "results"

try:
    from reporting.presentation_parser import parse_presentation
except ImportError:
    try:
        from presentation_parser import parse_presentation
    except ImportError:
        parse_presentation = None

logger = get_logger("presentation_generator")

# ---------------------------------------------------------------------------
# Paleta de cores
# ---------------------------------------------------------------------------
AZUL_ESCURO   = RGBColor(0x1e, 0x3a, 0x8a)  # #1e3a8a — fundo de seção
AZUL_MEDIO    = RGBColor(0x1d, 0x4e, 0xd8)  # #1d4ed8 — destaques
AZUL_CLARO    = RGBColor(0xdb, 0xea, 0xfe)  # #dbeafe — fundo suave
BRANCO        = RGBColor(0xff, 0xff, 0xff)
PRETO         = RGBColor(0x11, 0x18, 0x27)  # quase preto
CINZA_CLARO   = RGBColor(0xf1, 0xf5, 0xf9)
CINZA_TEXTO   = RGBColor(0x47, 0x55, 0x69)
VERDE_SOTA    = RGBColor(0x16, 0xa3, 0x4a)  # #16a34a — SOTA highlight
VERDE_CLARO   = RGBColor(0xdc, 0xfc, 0xe7)  # #dcfce7 — fundo linha SOTA
VERMELHO_ERRO = RGBColor(0xdc, 0x26, 0x26)  # #dc2626
AMARELO_AVISO = RGBColor(0xf5, 0x9e, 0x0b)

# ---------------------------------------------------------------------------
# Tamanho 16:9
# ---------------------------------------------------------------------------
W = Inches(13.33)
H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers básicos
# ---------------------------------------------------------------------------

def _rgb(color: RGBColor):
    return color

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # totalmente em branco
    return prs.slides.add_slide(blank_layout)

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

# ---------------------------------------------------------------------------
# Mode Filtering System — [modes: ...] tag support
# ---------------------------------------------------------------------------

def filter_markdown_by_mode(markdown_text: str, mode: str) -> str:
    """
    Filtra slides Markdown baseado em [modes: ...] tags.

    Sintaxe:
        [modes: full, compact]  → slide aparece em ambas
        [modes: full]           → slide só em full
        [modes: compact]        → slide só em compact

    Exemplo:
        [modes: full, compact]
        ## Slide Compartilhado
        Conteúdo aqui...

        [modes: full]
        ## Slide Só Full
        Apenas na versão completa...

    Args:
        markdown_text: conteúdo Markdown completo
        mode: "full" ou "compact"

    Returns:
        Markdown filtrado com tags [modes: ...] removidas
    """
    lines = markdown_text.split('\n')
    filtered_lines = []
    current_mode_match = None

    for line in lines:
        # Detectar tags [modes: ...]
        if line.strip().startswith('[modes:'):
            # Extrair modos: [modes: full, compact] ou [modes: full]
            modes_str = line.split('[modes:')[1].split(']')[0].strip()
            # Parse: "full, compact" → ["full", "compact"]
            modes_list = [m.strip() for m in modes_str.split(',')]
            # Remover labels se houver (ex: "full, compact | label: da_loss_full, da_loss_compact")
            modes_list = [m.split('|')[0].strip() for m in modes_list]
            current_mode_match = modes_list
            # Não incluir a linha da tag no output
            continue

        # Se há um modo definido, incluir apenas se for relevante
        if current_mode_match is not None:
            if mode in current_mode_match:
                filtered_lines.append(line)
            # Resetar quando encontrar próxima seção (---)
            if line.strip() == '---':
                current_mode_match = None
        else:
            # Sem tag [modes: ...], incluir por padrão
            filtered_lines.append(line)

    return '\n'.join(filtered_lines)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=PRETO, align=PP_ALIGN.LEFT,
                font_name="Calibri", wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font_name
    return txb

def add_para(tf, text, font_size=16, bold=False, color=PRETO,
             align=PP_ALIGN.LEFT, indent_level=0, font_name="Calibri",
             space_before=0, italic=False):
    """Adiciona parágrafo ao text_frame existente."""
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.level = indent_level
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = _Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font_name
    run.font.italic = italic
    return p

def add_title_bar(slide, title_text, subtitle_text=None,
                  bg_color=AZUL_ESCURO, text_color=BRANCO):
    """Barra de título azul no topo."""
    bar_h = Inches(1.3)
    add_rect(slide, 0, 0, W, bar_h, fill_color=bg_color)
    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), W - Inches(0.8), bar_h - Inches(0.1))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Calibri"
    if subtitle_text:
        add_para(tf, subtitle_text, font_size=16, color=RGBColor(0xba, 0xd3, 0xf8),
                 space_before=2)
    return txb

def add_code_box(slide, left, top, width, height, code_text, font_size=13):
    """Caixa de código monoespaçada com fundo cinza escuro."""
    add_rect(slide, left, top, width, height, fill_color=RGBColor(0x1e, 0x29, 0x3b))
    txb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1),
                                   width - Inches(0.3), height - Inches(0.2))
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code_text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)
    return txb

def add_table(slide, left, top, width, rows_data, col_widths=None,
              header_fill=AZUL_ESCURO, header_text_color=BRANCO,
              row_alt_fill=CINZA_CLARO, font_size=14,
              highlight_rows=None, highlight_color=VERDE_CLARO,
              highlight_text_color=VERDE_SOTA):
    """
    rows_data: list of lists of strings. Primeira linha = header.
    highlight_rows: set de índices de linha (0=header) a destacar em verde.
    """
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 1
    row_h  = Inches(0.42)
    total_h = row_h * n_rows

    if col_widths is None:
        col_w = width / n_cols
        col_widths = [col_w] * n_cols

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, total_h).table

    # Ajusta larguras das colunas
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = int(cw)

    for ri, row in enumerate(rows_data):
        is_header = (ri == 0)
        is_highlighted = highlight_rows and ri in highlight_rows
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            # Limpa texto negrito se houver **...**
            text_clean = cell_text.replace("**", "")
            has_bold = "**" in cell_text

            run = p.add_run()
            run.text = text_clean
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"

            if is_header:
                run.font.bold = True
                run.font.color.rgb = header_text_color
                _set_cell_fill(cell, header_fill)
            elif is_highlighted:
                run.font.bold = True
                run.font.color.rgb = highlight_text_color
                _set_cell_fill(cell, highlight_color)
            else:
                run.font.bold = has_bold
                run.font.color.rgb = PRETO
                if ri % 2 == 0:
                    _set_cell_fill(cell, row_alt_fill)
                else:
                    _set_cell_fill(cell, BRANCO)

    return tbl

def _set_cell_fill(cell, color: RGBColor):
    """Define cor de fundo de uma célula de tabela."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")

def footer(slide, text="FG2P · PT-BR G2P BiLSTM · 2026", page_num=None):
    """Rodapé discreto."""
    txt = text if page_num is None else f"{text}   |   {page_num}"
    add_textbox(slide, Inches(0.3), H - Inches(0.35), W - Inches(0.6), Inches(0.3),
                txt, font_size=10, color=RGBColor(0x94, 0xa3, 0xb8), align=PP_ALIGN.CENTER)


def _table_from_parsed(parsed_table: dict) -> list[list[str]]:
    """
    Converte formato de tabela do parser para o formato esperado por add_table().
    Parser: {"headers": [...], "rows": [[...], ...]}
    add_table: [headers, row1, row2, ...] (lista de listas)
    """
    if not parsed_table:
        return []
    return [parsed_table["headers"]] + parsed_table.get("rows", [])


# ---------------------------------------------------------------------------
# Slides individuais
# ---------------------------------------------------------------------------

def slide_opening(prs: Presentation, slide_data=None):
    """Slide de Abertura — Página de Título com Autores e Contexto.

    Se slide_data (dict com meta do parser) for passado, usa seus valores.
    Caso contrário, usa valores hardcoded como fallback.
    """
    meta = slide_data or {}
    author      = meta.get("author",      "Leonardo Marques de Souza")
    professor   = meta.get("professor",   "Prof. Dr. André Eugênio Lazzaretti")
    institution = meta.get("institution", "Universidade Tecnológica Federal do Paraná (UTFPR)")
    course      = meta.get("course",      "Deep Learning")
    year        = meta.get("year",        "Final de 2025")

    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=AZUL_ESCURO)

    add_textbox(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
                "FG2P — Conversão Grafema-para-Fonema",
                font_size=48, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.6),
                "Português Brasileiro com BiLSTM + Distance-Aware Loss",
                font_size=28, color=RGBColor(0x93, 0xc5, 0xfd), align=PP_ALIGN.CENTER)

    add_rect(s, Inches(3), Inches(3.0), Inches(7.33), Inches(0.05),
             fill_color=RGBColor(0x60, 0xa5, 0xfa))

    add_textbox(s, Inches(1.0), Inches(3.8), Inches(11.33), Inches(0.5),
                f"Autor: {author}",
                font_size=24, color=BRANCO, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1.0), Inches(4.4), Inches(11.33), Inches(0.5),
                professor,
                font_size=22, color=RGBColor(0x93, 0xc5, 0xfd), align=PP_ALIGN.CENTER)

    add_rect(s, Inches(2.0), Inches(5.3), Inches(9.33), Inches(1.0),
             fill_color=RGBColor(0x1d, 0x4e, 0xd8), line_color=RGBColor(0x60, 0xa5, 0xfa))

    add_textbox(s, Inches(2.1), Inches(5.45), Inches(9.13), Inches(0.8),
                f"Atividade Acadêmica: {course}\n{institution}\n{year}",
                font_size=18, color=BRANCO, align=PP_ALIGN.CENTER)

    footer(s)
    return s


def slide_title(prs: Presentation):
    """Slide 1 — Título"""
    s = blank_slide(prs)
    # Fundo gradiente simulado: dois retângulos
    add_rect(s, 0, 0, W, H, fill_color=AZUL_ESCURO)
    add_rect(s, 0, H - Inches(2.5), W, Inches(2.5), fill_color=RGBColor(0x17, 0x30, 0x78))

    # Linha decorativa
    add_rect(s, Inches(0.5), Inches(2.6), Inches(6), Inches(0.06),
             fill_color=RGBColor(0x60, 0xa5, 0xfa))

    add_textbox(s, Inches(0.5), Inches(1.0), Inches(12), Inches(1.0),
                "FG2P", font_size=60, bold=True, color=BRANCO,
                align=PP_ALIGN.LEFT, font_name="Calibri")

    add_textbox(s, Inches(0.5), Inches(2.0), Inches(12), Inches(0.7),
                "Conversão Grafema-Fonema para o Português Brasileiro",
                font_size=26, color=RGBColor(0xba, 0xd3, 0xf8),
                align=PP_ALIGN.LEFT)

    add_textbox(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                "Modelo BiLSTM Encoder-Decoder com Distance-Aware Loss",
                font_size=18, color=RGBColor(0x93, 0xc5, 0xfd), align=PP_ALIGN.LEFT)

    # Caixa de resultado em destaque
    add_rect(s, Inches(0.5), Inches(3.9), Inches(7.5), Inches(0.9),
             fill_color=RGBColor(0x1d, 0x4e, 0xd8))
    add_textbox(s, Inches(0.6), Inches(3.95), Inches(7.3), Inches(0.8),
                '"computador"  →  k õ p u t a ˈ d o x',
                font_size=20, bold=True, color=BRANCO,
                font_name="Courier New", align=PP_ALIGN.LEFT)

    # Stats
    txb = s.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "PER 0,49%  ·  WER 4,96%  ·  28.782 palavras de teste  ·  100% em OOV PT-BR"
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)
    r.font.name = "Calibri"

    return s


def slide_g2p_problem(prs: Presentation, num: int):
    """Slide 2 — O que é G2P?"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "O que é o problema G2P?", "Grafema → Fonema")

    # Dois exemplos lado a lado
    for i, (entrada, saida, nota) in enumerate([
        ("casa", "k  a  z  a", "c antes de vogal posterior → /k/"),
        ("cena", "s  e  n  ə", "c antes de vogal anterior → /s/"),
    ]):
        col_left = Inches(0.5) + i * Inches(6.0)
        add_rect(s, col_left, Inches(1.6), Inches(5.6), Inches(2.5),
                 fill_color=AZUL_CLARO)
        add_textbox(s, col_left + Inches(0.2), Inches(1.7), Inches(5.2), Inches(0.5),
                    f'Entrada:  "{entrada}"', font_size=18, bold=True, color=AZUL_ESCURO)
        add_code_box(s, col_left + Inches(0.2), Inches(2.25), Inches(5.2), Inches(0.65),
                     saida, font_size=20)
        add_textbox(s, col_left + Inches(0.2), Inches(3.0), Inches(5.2), Inches(0.5),
                    nota, font_size=14, color=CINZA_TEXTO, italic=True)

    # Conclusão
    add_rect(s, Inches(0.5), Inches(4.3), W - Inches(1.0), Inches(0.9),
             fill_color=AZUL_ESCURO)
    add_textbox(s, Inches(0.7), Inches(4.4), W - Inches(1.4), Inches(0.7),
                "O modelo não memoriza palavras — aprende regras que generalizam.",
                font_size=19, bold=True, color=BRANCO)

    txb2 = s.shapes.add_textbox(Inches(0.5), Inches(5.4), W - Inches(1.0), Inches(1.5))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    add_para(tf2, "Grafema: unidade escrita (letra ou dígrafo) · Fonema: unidade sonora da língua",
             font_size=14, color=CINZA_TEXTO)
    add_para(tf2, "IPA: Alfabeto Fonético Internacional — notação padrão usada neste trabalho",
             font_size=14, color=CINZA_TEXTO, space_before=4)

    footer(s, page_num=f"2")
    return s


def slide_ptbr_hard(prs: Presentation, num: int, slide_data=None):
    """Slide 3 — Por que PT-BR é difícil?"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Por que PT-BR é difícil?", "3 desafios principais")

    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows = _table_from_parsed(slide_data["tables"][0])
    else:
        rows = [
            ["Desafio", "Exemplo", "Regra"],
            ["Ambiguidade grafêmica", '"cama" vs "cena"', "c → /k/ ou /s/"],
            ["Dependência de posição", '"roda" vs "carro"', "r inicial→x, rr→x, r medial→ɾ"],
            ["Neutralização vocálica", '"pé" vs "pedra"', "ɛ tônico → e/ɪ átono"],
        ]
    col_widths = [Inches(4.0), Inches(4.0), Inches(4.6)]
    add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), rows,
              col_widths=col_widths, font_size=16)

    add_rect(s, Inches(0.5), Inches(4.1), W - Inches(1.0), Inches(1.0),
             fill_color=RGBColor(0xfe, 0xf0, 0x8a))
    add_textbox(s, Inches(0.7), Inches(4.15), W - Inches(1.4), Inches(0.9),
                "~60% dos erros do modelo são neutralizações vocálicas legítimas"
                " — ambiguidade irredutível sem contexto prosódico",
                font_size=17, bold=True, color=RGBColor(0x78, 0x35, 0x00))

    txb = s.shapes.add_textbox(Inches(0.5), Inches(5.3), W - Inches(1.0), Inches(1.5))
    tf = txb.text_frame
    tf.word_wrap = True
    add_para(tf, "Neutralização vocálica: vogal muda de timbre dependendo da posição na palavra",
             font_size=14, color=CINZA_TEXTO)
    add_para(tf, "Palatalização: t/d→tʃ/dʒ antes de i, como em 'tia' → /tʃ i a/",
             font_size=14, color=CINZA_TEXTO, space_before=4)

    footer(s, page_num=str(num))
    return s


def slide_data(prs: Presentation, num: int, slide_data=None):
    """Slide 4 — Os Dados"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Os Dados")

    # Painel esquerdo
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.8), Inches(4.8),
             fill_color=AZUL_CLARO)
    txb = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(5.4), Inches(4.5))
    tf = txb.text_frame
    tf.word_wrap = True
    add_para(tf, "95.937 pares (palavra, IPA)", font_size=20, bold=True,
             color=AZUL_ESCURO)
    add_para(tf, "Fonte: dicionário fonético PT-BR", font_size=15, color=CINZA_TEXTO,
             space_before=4)
    add_para(tf, "Charset: a–z (exceto k,w,y) + ç á à â ã é ê í ó ô õ ú ü",
             font_size=14, color=CINZA_TEXTO, space_before=6)
    add_para(tf, "Divisão estratificada 60/10/30", font_size=16, bold=True,
             color=AZUL_ESCURO, space_before=12)

    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows_div = _table_from_parsed(slide_data["tables"][0])
    else:
        rows_div = [
            ["Subconjunto", "Tamanho"],
            ["Treino", "57.561"],
            ["Validação", "9.594"],
            ["Teste", "28.782"],
        ]

    add_table(s, Inches(0.6), Inches(3.8), Inches(5.0), rows_div,
              col_widths=[Inches(3.0), Inches(2.0)], font_size=15,
              highlight_rows={3}, highlight_color=AZUL_CLARO,
              highlight_text_color=AZUL_ESCURO)

    # Painel direito — descoberta
    add_rect(s, Inches(6.6), Inches(1.5), Inches(6.3), Inches(2.2),
             fill_color=VERDE_CLARO)
    add_textbox(s, Inches(6.8), Inches(1.55), Inches(5.9), Inches(0.5),
                "Descoberta inesperada:", font_size=17, bold=True,
                color=VERDE_SOTA)
    add_textbox(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(1.5),
                "Split 60/10/30 supera 70/10/20 em −41% PER\n"
                "→ Mais dados de treino ≠ modelo melhor\n"
                "→ Teste maior = medição estatística melhor",
                font_size=15, color=PRETO)

    add_rect(s, Inches(6.6), Inches(4.0), Inches(6.3), Inches(2.0),
             fill_color=CINZA_CLARO)
    add_textbox(s, Inches(6.8), Inches(4.05), Inches(5.9), Inches(0.4),
                "Qualidade do split:", font_size=15, bold=True, color=AZUL_ESCURO)
    add_textbox(s, Inches(6.8), Inches(4.5), Inches(5.9), Inches(1.4),
                "χ² = 0,95  (p = 0,678)\nCramér V = 0,0007\n→ distribuição fonológica balanceada",
                font_size=15, color=CINZA_TEXTO)

    footer(s, page_num=str(num))
    return s


def slide_architecture(prs: Presentation, num: int):
    """Slide 5 — A Arquitetura"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "A Arquitetura", "BiLSTM Encoder-Decoder + Atenção Bahdanau")

    code = (
        '"c a s a"   ← grafemas de entrada\n'
        '    ↓\n'
        '[ Embedding ]     192D: cada letra → vetor\n'
        '    ↓\n'
        '[ BiLSTM Encoder ]  2 camadas, 384D\n'
        '  → lê a palavra nos dois sentidos\n'
        '  → "c" em "cama" ≠ "c" em "cena"\n'
        '    ↓\n'
        '[ Atenção Bahdanau ]\n'
        '  → decoder "olha" para onde precisa\n'
        '  → aprende: "nh" → foco nas duas letras\n'
        '    ↓\n'
        '[ LSTM Decoder ]  gera fonema por fonema\n'
        '    ↓\n'
        'k  a  z  a   ← fonemas de saída'
    )
    add_code_box(s, Inches(0.4), Inches(1.4), Inches(6.4), Inches(5.6),
                 code, font_size=13)

    # Anotações à direita
    bullets = [
        ("BiLSTM", "Lê a palavra da esquerda e da direita — o contexto de 'c' em 'cena' inclui a vogal que vem depois"),
        ("Encoder-Decoder", "Encoder = compreende a entrada · Decoder = gera a saída passo a passo"),
        ("Atenção Bahdanau", "Mecanismo de consulta: para cada fonema gerado, o decoder pondera toda a sequência de entrada"),
        ("9,7M parâmetros", "Modelo intermediário — 3× maior que baseline, sem overfitting"),
    ]
    y_pos = Inches(1.4)
    for title, desc in bullets:
        add_rect(s, Inches(7.1), y_pos, Inches(5.8), Inches(1.1),
                 fill_color=AZUL_CLARO)
        add_textbox(s, Inches(7.25), y_pos + Inches(0.05), Inches(5.5), Inches(0.4),
                    title, font_size=15, bold=True, color=AZUL_ESCURO)
        add_textbox(s, Inches(7.25), y_pos + Inches(0.45), Inches(5.5), Inches(0.6),
                    desc, font_size=13, color=CINZA_TEXTO)
        y_pos += Inches(1.2)

    footer(s, page_num=str(num))
    return s


def slide_attention(prs: Presentation, num: int):
    """Slide 6 — Por que Atenção?"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Por que Atenção?")

    # Sem atenção
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.8), Inches(2.2),
             fill_color=RGBColor(0xfe, 0xe2, 0xe2))
    add_textbox(s, Inches(0.6), Inches(1.55), Inches(5.5), Inches(0.4),
                "Sem atenção:", font_size=17, bold=True,
                color=VERMELHO_ERRO)
    add_code_box(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(1.4),
                 '"biscoito" → encoder comprime em 1 vetor (384D)\nDecoder não consegue "ir atrás"\npara verificar "sc" → padrão /sk/',
                 font_size=13)

    # Com atenção
    add_rect(s, Inches(0.4), Inches(3.9), Inches(5.8), Inches(2.8),
             fill_color=VERDE_CLARO)
    add_textbox(s, Inches(0.6), Inches(3.95), Inches(5.5), Inches(0.4),
                "Com atenção Bahdanau:", font_size=17, bold=True, color=VERDE_SOTA)
    add_code_box(s, Inches(0.6), Inches(4.4), Inches(5.5), Inches(1.7),
                 'Ao gerar /k/:\n'
                 '  s=0.70 (olha em s)\n'
                 '  c=0.20 (olha em c)\n'
                 '  Aprende: "sc" = padrão /sk/',
                 font_size=13)

    # Exemplos lado direito
    txb = s.shapes.add_textbox(Inches(6.6), Inches(1.5), Inches(6.3), Inches(5.0))
    tf = txb.text_frame
    tf.word_wrap = True
    add_para(tf, "Essencial para padrões n:1", font_size=18, bold=True,
             color=AZUL_ESCURO)
    add_para(tf, "", font_size=8)

    exemplos = [
        ("ch → /ʃ/", "2 grafemas → 1 fonema"),
        ("nh → /ɲ/", "2 grafemas → 1 fonema"),
        ("lh → /ʎ/", "2 grafemas → 1 fonema"),
        ("x", "/ʃ/ ou /ks/ dependendo do contexto"),
        ("r / rr", "posição determina o rótico"),
    ]
    for gr, desc in exemplos:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        r1 = p.add_run()
        r1.text = f"  {gr}:  "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = AZUL_ESCURO
        r1.font.name = "Courier New"
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = CINZA_TEXTO
        r2.font.name = "Calibri"

    footer(s, page_num=str(num))
    return s


def slide_ce_problem(prs: Presentation, num: int):
    """Slide 7 — O Problema da CrossEntropy"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "O Problema da CrossEntropy",
                  "Todos os erros são tratados como igualmente ruins")

    add_code_box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.5),
                 'Cenário A — Quase certo:\n'
                 '  ɛ → e  (mesma posição, só altura)  penalidade = 1.0\n\n'
                 'Cenário B — Completamente errado:\n'
                 '  k → a  (oclusiva vs vogal)  penalidade = 1.0  ← MESMO VALOR!',
                 font_size=16)

    add_rect(s, Inches(0.5), Inches(4.2), W - Inches(1.0), Inches(0.9),
             fill_color=VERMELHO_ERRO)
    add_textbox(s, Inches(0.7), Inches(4.27), W - Inches(1.4), Inches(0.75),
                "CrossEntropy não sabe que certos erros são 'quase corretos'"
                " — não diferencia proximidade fonológica",
                font_size=18, bold=True, color=BRANCO)

    txb = s.shapes.add_textbox(Inches(0.5), Inches(5.3), W - Inches(1.0), Inches(1.7))
    tf = txb.text_frame
    tf.word_wrap = True
    add_para(tf, "CrossEntropy: L = −log(p(y_correto))  —  mede apenas se o correto foi escolhido",
             font_size=15, color=CINZA_TEXTO)
    add_para(tf, "Problema: a penalidade é binária (0 se correto, 1 se errado) — não reflete distância articulatória",
             font_size=15, color=CINZA_TEXTO, space_before=6)
    add_para(tf, "Consequência: o modelo não aprende a 'preferir erros inteligentes' quando está errado",
             font_size=15, color=CINZA_TEXTO, space_before=6)

    footer(s, page_num=str(num))
    return s


def slide_da_loss_problem(prs: Presentation, num: int):
    """Slide — Distance-Aware Loss: Problema com CrossEntropy"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Distance-Aware Loss — Problema",
                  "Por que CrossEntropy penaliza erros iguais?")

    # Problema
    add_rect(s, Inches(0.5), Inches(1.5), W - Inches(1.0), Inches(0.8),
             fill_color=AZUL_CLARO)
    add_textbox(s, Inches(0.7), Inches(1.55), W - Inches(1.4), Inches(0.7),
                "CrossEntropy trata todos os erros igualmente",
                font_size=18, bold=True, color=AZUL_ESCURO)

    # Exemplos de erros
    txb = s.shapes.add_textbox(Inches(0.7), Inches(2.5), W - Inches(1.4), Inches(3.0))
    tf = txb.text_frame
    tf.word_wrap = True

    add_para(tf, "Erros do LSTM no PT-BR:", font_size=16, bold=True, color=AZUL_ESCURO)
    add_para(tf, "e → ɛ :  Pequeno — confusão vocálica natural (mesma altura, mesma posição)",
             font_size=14, color=CINZA_TEXTO)
    add_para(tf, "e → x :  GRANDE — sons completamente diferentes (vogal vs fricativa velar)",
             font_size=14, color=CINZA_TEXTO, space_before=12)

    add_para(tf, "Impacto quantificado:", font_size=16, bold=True, color=AZUL_ESCURO,
             space_before=12)
    add_para(tf, "462 erros e↔ɛ = 18.9% de TODOS os erros do Exp104b",
             font_size=14, color=VERMELHO_ERRO, bold=True, font_name="Courier New")

    add_para(tf, "Solução: ponderar loss pela distância articulatória",
             font_size=15, bold=True, color=VERDE_SOTA, space_before=12)

    footer(s, page_num=str(num))
    return s


def slide_da_loss(prs: Presentation, num: int, slide_data=None):
    """Slide 8 — A Distance-Aware Loss: Fórmula e Componentes"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Distance-Aware Loss — Fórmula",
                  "Componentes e distâncias articulatórias")

    # Fórmula em destaque
    add_rect(s, Inches(0.5), Inches(1.5), W - Inches(1.0), Inches(0.9),
             fill_color=AZUL_ESCURO)
    add_textbox(s, Inches(0.7), Inches(1.55), W - Inches(1.4), Inches(0.8),
                "L  =  L_CE  +  λ · d(ŷ, y) · p(ŷ)",
                font_size=28, bold=True, color=BRANCO,
                font_name="Courier New", align=PP_ALIGN.CENTER)

    # Tabela de componentes
    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows = _table_from_parsed(slide_data["tables"][0])
    else:
        rows = [
            ["Componente", "Significado", "Escala"],
            ["L_CE", "Penalidade base (CrossEntropy)", "0–∞"],
            ["d(ŷ, y)", "Distância articulatória PanPhon entre predito e correto", "0–1"],
            ["p(ŷ)", "Confiança do modelo no token predito (softmax)", "0–1"],
            ["λ = 0,20", "Peso do sinal fonológico — sweet spot empírico", "hiperparâmetro"],
        ]

    col_widths = [Inches(2.5), Inches(7.0), Inches(2.8)]
    add_table(s, Inches(0.5), Inches(2.55), Inches(12.3), rows,
              col_widths=col_widths, font_size=15,
              highlight_rows={4}, highlight_color=VERDE_CLARO,
              highlight_text_color=VERDE_SOTA)

    # Distâncias reais
    txb = s.shapes.add_textbox(Inches(0.5), Inches(5.35), W - Inches(1.0), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = True
    add_para(tf, "Distâncias reais PanPhon (exemplos):", font_size=15, bold=True, color=AZUL_ESCURO)
    add_para(tf, "e ↔ ɛ ≈ 0.10   p ↔ b ≈ 0.05   s ↔ ʃ ≈ 0.15   a ↔ k ≈ 0.90   . ↔ ˈ = 0.00 (corrigido no Exp104b)",
             font_size=14, color=CINZA_TEXTO, space_before=4,
             font_name="Courier New" if False else "Calibri")

    footer(s, page_num=str(num))
    return s


def slide_da_example(prs: Presentation, num: int):
    """Slide 9 — DA Loss passo a passo"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "DA Loss — Exemplo Passo a Passo",
                  'Palavra: "cena" · Posição: grafema "c" · Correto: /s/')

    code = (
        'PASSO 1 — Modelo produz distribuição de probabilidade:\n'
        '  s=42%   ʃ=35%   z=12%   t=6%   ...\n'
        '  ↓ argmax\n'
        '  Predição: ʃ  (errado!)\n'
        '\n'
        'PASSO 2 — CrossEntropy (olha só para o correto):\n'
        '  L_CE = −log(0.42) = 0.87\n'
        '\n'
        'PASSO 3 — Distância articulatória PanPhon:\n'
        '  d(ʃ, s) = 0.15   ← fricativas, mesmo modo, ponto diferente\n'
        '\n'
        'PASSO 4 — Termo extra de DA Loss:\n'
        '  0.20 × 0.15 × 0.35 = 0.010\n'
        '\n'
        'TOTAL:  L = 0.87 + 0.01 = 0.88'
    )
    add_code_box(s, Inches(0.4), Inches(1.4), Inches(7.0), Inches(5.6),
                 code, font_size=13)

    # Comparação à direita
    add_rect(s, Inches(7.6), Inches(1.4), Inches(5.3), Inches(2.5),
             fill_color=VERDE_CLARO)
    add_textbox(s, Inches(7.8), Inches(1.45), Inches(5.0), Inches(0.4),
                "Se tivesse predito /k/:", font_size=16, bold=True,
                color=VERDE_SOTA)
    add_code_box(s, Inches(7.8), Inches(1.9), Inches(5.0), Inches(0.85),
                 'd(k, s) = 0.80  →  extra = 0.056\nTotal = 0.87 + 0.056 = 0.93',
                 font_size=14)

    add_rect(s, Inches(7.6), Inches(4.0), Inches(5.3), Inches(2.8),
             fill_color=AZUL_CLARO)
    add_textbox(s, Inches(7.8), Inches(4.05), Inches(5.0), Inches(0.4),
                "O que isso significa:", font_size=16, bold=True, color=AZUL_ESCURO)
    txb = s.shapes.add_textbox(Inches(7.8), Inches(4.5), Inches(5.0), Inches(2.0))
    tf = txb.text_frame
    tf.word_wrap = True
    add_para(tf, "5% de diferença por passo", font_size=15, bold=True, color=AZUL_ESCURO)
    add_para(tf, "Após milhares de batches: o modelo aprende a 'desempatar para o lado fonologicamente certo'",
             font_size=14, color=CINZA_TEXTO, space_before=4)
    add_para(tf, "ʃ → s (5% de penalidade extra)  =  leve correção", font_size=13,
             color=VERDE_SOTA, space_before=8)
    add_para(tf, "k → s (50% de penalidade extra) =  forte correção", font_size=13,
             color=VERMELHO_ERRO, space_before=4)

    footer(s, page_num=str(num))
    return s


def slide_confidence_factor(prs: Presentation, num: int):
    """Slide 10 — Por que multiplicar por p(ŷ)?"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Por que multiplicar por p(ŷ)?",
                  "O fator de confiança evita punir incerteza honesta")

    add_code_box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.8),
                 'Época 5 — modelo incerto:\n'
                 '  Prediz /k/ com p=52%  →  extra = λ · d · 0.52  →  penalidade moderada\n'
                 '  "Ainda aprendendo — só corrigindo a direção"\n'
                 '\n'
                 'Época 40 — modelo confiante:\n'
                 '  Prediz /k/ com p=91%  →  extra = λ · d · 0.91  →  penalidade alta\n'
                 '  "Você está MUITO certo de algo completamente errado"',
                 font_size=16)

    add_rect(s, Inches(0.5), Inches(4.5), W - Inches(1.0), Inches(0.85),
             fill_color=AZUL_CLARO)
    add_textbox(s, Inches(0.7), Inches(4.55), W - Inches(1.4), Inches(0.75),
                "Análogo ao Label Smoothing, mas proporcional à distância fonológica — não uniforme",
                font_size=18, bold=True, color=AZUL_ESCURO)

    txb = s.shapes.add_textbox(Inches(0.5), Inches(5.5), W - Inches(1.0), Inches(1.5))
    tf = txb.text_frame
    tf.word_wrap = True
    add_para(tf, "Label Smoothing: distribui um ε fixo entre todos os tokens — não distingue k de ʃ como candidatos",
             font_size=14, color=CINZA_TEXTO)
    add_para(tf, "DA Loss: o 'suavizamento' é proporcional a d(ŷ,y) — erros articulatoriamente distantes recebem mais correção",
             font_size=14, color=CINZA_TEXTO, space_before=6)

    footer(s, page_num=str(num))
    return s


def slide_separators(prs: Presentation, num: int, slide_data=None):
    """Slide 11 — Separadores Silábicos"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Separadores Silábicos — O Trade-off",
                  "Adicionando '.' como token de fronteira silábica")

    add_code_box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2),
                 'Exp9   (sem sep):  k  õ  p  u  t  a  ˈ  d  o  x\n'
                 'Exp102 (com sep):  k  .  õ p  .  u  .  t a  .  ˈ  d  o  x  .',
                 font_size=17)

    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows = _table_from_parsed(slide_data["tables"][0])
    else:
        rows = [
            ["Experimento", "PER ↓", "WER ↓", "Comentário"],
            ["Exp9  (sem sep)", "0,58%", "4,96%", "SOTA WER ← menos erros de palavra inteira"],
            ["Exp102 (com sep)", "0,52%", "5,79%", "SOTA PER ← menos erros de fonema"],
            ["Exp103 (DA+sep)", "0,53%", "5,73%", "Efeitos não-aditivos — trade-off permanece"],
        ]

    col_widths = [Inches(3.5), Inches(1.8), Inches(1.8), Inches(5.2)]
    add_table(s, Inches(0.5), Inches(2.85), Inches(12.3), rows,
              col_widths=col_widths, font_size=15,
              highlight_rows={1}, highlight_color=VERDE_CLARO,
              highlight_text_color=VERDE_SOTA)

    add_rect(s, Inches(0.5), Inches(4.85), W - Inches(1.0), Inches(0.85),
             fill_color=AZUL_ESCURO)
    add_textbox(s, Inches(0.7), Inches(4.9), W - Inches(1.4), Inches(0.75),
                "Cada '.' mal-posicionado → palavra inteira errada (WER binário)"
                "   →   Trade-off Pareto fundamental",
                font_size=17, bold=True, color=BRANCO)

    txb = s.shapes.add_textbox(Inches(0.5), Inches(5.9), W - Inches(1.0), Inches(1.2))
    tf = txb.text_frame
    tf.word_wrap = True
    add_para(tf, "Separador silábico: token '.' inserido entre sílabas no alvo de treino (ex: com.pu.ta.dor)",
             font_size=14, color=CINZA_TEXTO)
    add_para(tf, "PER conta erros em fonemas · WER conta palavras com qualquer erro — um '.' extra = WER+=1",
             font_size=14, color=CINZA_TEXTO, space_before=4)

    footer(s, page_num=str(num))
    return s


def slide_custom_dist(prs: Presentation, num: int, slide_data=None):
    """Slide 12 — Distâncias Customizadas"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Distâncias Customizadas — O Bug que Virou Feature")

    # Problema
    add_rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(2.4),
             fill_color=RGBColor(0xfe, 0xe2, 0xe2))
    add_textbox(s, Inches(0.6), Inches(1.45), Inches(5.5), Inches(0.4),
                "Problema: '.' e 'ˈ' têm vetor ZERO no PanPhon",
                font_size=16, bold=True, color=VERMELHO_ERRO)
    add_code_box(s, Inches(0.6), Inches(1.9), Inches(5.5), Inches(1.5),
                 'd(., ˈ) = 0.0  ← Loss NÃO penaliza confusão!\n'
                 'd(., a) = 0.0  ← indistinguíveis da Loss',
                 font_size=14)

    # Solução
    add_rect(s, Inches(6.6), Inches(1.4), Inches(6.3), Inches(2.4),
             fill_color=VERDE_CLARO)
    add_textbox(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(0.4),
                "Solução (Exp104b): override pós-normalização",
                font_size=16, bold=True, color=VERDE_SOTA)
    add_code_box(s, Inches(6.8), Inches(1.9), Inches(6.0), Inches(1.5),
                 'for sym in {".", "ˈ"}:\n'
                 '    for other in vocab:\n'
                 '        dist[sym][other] = 1.0  # máximo',
                 font_size=14)

    # Tabela de resultados
    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows = _table_from_parsed(slide_data["tables"][0])
    else:
        rows = [
            ["Experimento", "PER", "WER", "Erros .↔ˈ"],
            ["Exp103 (sem override)", "0,53%", "5,73%", "~107"],
            ["Exp104 (bug: override antes norm.)", "0,54%", "5,88%", "~119"],
            ["**Exp104b (override correto)**", "**0,49%**", "**5,43%**", "~106"],
        ]

    col_widths = [Inches(5.5), Inches(1.8), Inches(1.8), Inches(2.5)]
    add_table(s, Inches(0.5), Inches(4.0), Inches(11.6), rows,
              col_widths=col_widths, font_size=15,
              highlight_rows={3}, highlight_color=VERDE_CLARO,
              highlight_text_color=VERDE_SOTA)

    add_textbox(s, Inches(0.5), Inches(6.1), W - Inches(1.0), Inches(0.5),
                "→ NOVO SOTA PER: 0,49%  (mesmo sem eliminar as confusões estruturais)",
                font_size=16, bold=True, color=VERDE_SOTA)

    footer(s, page_num=str(num))
    return s


def slide_evolution(prs: Presentation, num: int):
    """Slide 13 — Evolução dos experimentos"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Evolução dos Experimentos", "PER ao longo das fases")

    add_code_box(s, Inches(0.5), Inches(1.45), Inches(7.5), Inches(5.55),
                 'PER (%)\n'
                 '1.12 │ ● Exp0 (70/10/20 — split errado)\n'
                 '     │\n'
                 '0.66 │         ● Exp1 (60/10/30)  −41%\n'
                 '0.63 │                 ● Exp5 (9.7M params)\n'
                 '0.61 │                         ● Exp7 (λ=0.20)\n'
                 '0.60 │                 ● Exp2 (17.2M params)\n'
                 '0.58 │                                 ● Exp9   ← SOTA WER 4.96%\n'
                 '0.52 │                                         ● Exp102 (sep)\n'
                 '0.49 │                                                 ● Exp104b ← SOTA PER\n'
                 '     └────────────────────────────────────────────────────────\n'
                 '      Fase1   Fase2   Fase3   Fase4   Fase5   Fase6',
                 font_size=13)

    # Insights à direita
    txb = s.shapes.add_textbox(Inches(8.3), Inches(1.45), Inches(4.7), Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True

    insights = [
        ("Fase 1", "Split 60/10/30 resolve problema de medição"),
        ("Fase 2-3", "Scaling não ajuda; DA Loss sim"),
        ("Fase 4", "λ=0,20 é sweet spot empírico"),
        ("Fase 5", "DA Loss no modelo 9.7M → SOTA WER"),
        ("Fase 6A", "Separadores → trade-off Pareto"),
        ("Fase 6B", "Custom dist fix → SOTA PER 0,49%"),
    ]
    for fase, desc in insights:
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        r1 = p.add_run()
        r1.text = f"{fase}: "
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = AZUL_ESCURO
        r1.font.name = "Calibri"
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = CINZA_TEXTO
        r2.font.name = "Calibri"

    footer(s, page_num=str(num))
    return s


def slide_ranking(prs: Presentation, num: int, slide_data=None):
    """Slide 14 — Ranking Final"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Ranking Final", "Qual modelo usar e quando")

    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows = _table_from_parsed(slide_data["tables"][0])
    else:
        rows = [
            ["Rank", "Modelo", "Params", "PER", "WER", "Speed (w/s)", "Para usar quando..."],
            ["**PER SOTA**", "Exp104b", "9,7M", "**0,49%**", "5,43%", "32,4", "TTS, alinhamento fonético, análise"],
            ["**WER SOTA**", "Exp9",    "9,7M", "0,58%", "**4,96%**", "40,3", "NLP, busca fonética, indexação"],
            ["Budget",       "Exp6",    "4,3M", "0,63%", "5,35%", "40,3", "Recursos computacionais limitados"],
            ["Baseline",     "Exp1",    "4,3M", "0,66%", "5,65%", "25,5", "Comparação simples"],
            ["Evitar",       "Exp0",    "4,3M", "1,12%", "9,37%", "23,9", "Split 70/10/20 — medição inadequada"],
        ]

    col_widths = [Inches(1.4), Inches(1.6), Inches(1.3), Inches(1.0), Inches(1.0), Inches(1.2), Inches(4.5)]
    add_table(s, Inches(0.4), Inches(1.45), Inches(12.5), rows,
              col_widths=col_widths, font_size=14,
              highlight_rows={1, 2})

    # CLI
    add_rect(s, Inches(0.4), Inches(5.5), W - Inches(0.8), Inches(0.75),
             fill_color=RGBColor(0x1e, 0x29, 0x3b))
    add_textbox(s, Inches(0.6), Inches(5.55), W - Inches(1.2), Inches(0.65),
                "python src/inference_light.py --index 18 --word computador   "
                "|   --index 11 para WER SOTA",
                font_size=15, color=RGBColor(0xe2, 0xe8, 0xf0),
                font_name="Courier New")

    footer(s, page_num=str(num))
    return s


def slide_generalization_design(prs: Presentation, num: int, slide_data=None):
    """Slide 15 — Design da avaliação de generalização"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Avaliação de Generalização — Design",
                  "O modelo memorizou ou aprendeu regras?")

    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows = _table_from_parsed(slide_data["tables"][0])
    else:
        rows = [
            ["Categoria", "N", "Objetivo"],
            ["Generalização PT-BR", "9", "Neologismos e portmanteaux — chars no vocab"],
            ["Consoantes duplas", "5", "lazzaretti, cappuccino → redução de geminadas"],
            ["Anglicismos (chars OK)", "5", "mouse, site → fonologia inglesa parcial"],
            ["Chars OOV (k/w/y)", "3", "wifi, yoga → FALHA ESPERADA e documentada"],
            ["PT-BR reais (OOV)", "5", "puxadinho, zunido → PROVA DE GENERALIZAÇÃO"],
            ["Controles", "4", "biscoito, computador → verificação de sanidade"],
        ]

    col_widths = [Inches(4.5), Inches(0.8), Inches(7.0)]
    add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), rows,
              col_widths=col_widths, font_size=15,
              highlight_rows={5})

    add_rect(s, Inches(0.5), Inches(5.75), W - Inches(1.0), Inches(0.75),
             fill_color=AZUL_CLARO)
    add_textbox(s, Inches(0.7), Inches(5.80), W - Inches(1.4), Inches(0.65),
                "31 palavras curadas · 6 categorias · avaliação qualitativa + score fonológico",
                font_size=16, bold=False, color=AZUL_ESCURO)

    footer(s, page_num=str(num))
    return s


def slide_oov_result(prs: Presentation, num: int, slide_data=None):
    """Slide 16 — O resultado mais importante"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "O Resultado que Mais Importa",
                  "Palavras PT-BR reais fora do vocabulário de treino: 5/5")

    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows = _table_from_parsed(slide_data["tables"][0])
    else:
        rows = [
            ["Palavra", "Predição do modelo", "Esperado", ""],
            ["puxadinho", "p u ʃ a ˈ d ʒ ĩ ɲ ʊ", "p u ʃ a ˈ d ʒ ĩ ɲ ʊ", "✓"],
            ["malcriado",  "m a w k ɾ i ˈ a d ʊ", "m a w k ɾ i ˈ a d ʊ", "✓"],
            ["arrombado",  "a x õ ˈ b a d ʊ",     "a x õ ˈ b a d ʊ",     "✓"],
            ["abacatada",  "a b a k a ˈ t a d ə", "a b a k a ˈ t a d ə", "✓"],
            ["zunido",     "z u ˈ n i d ʊ",        "z u ˈ n i d ʊ",       "✓"],
        ]

    col_widths = [Inches(2.5), Inches(4.5), Inches(4.5), Inches(0.8)]
    add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), rows,
              col_widths=col_widths, font_size=15,
              highlight_rows={1, 2, 3, 4, 5},
              highlight_color=VERDE_CLARO, highlight_text_color=VERDE_SOTA)

    add_rect(s, Inches(0.5), Inches(5.3), W - Inches(1.0), Inches(0.85),
             fill_color=VERDE_SOTA)
    add_textbox(s, Inches(0.7), Inches(5.37), W - Inches(1.4), Inches(0.75),
                "O modelo acerta: x→ʃ, palatalização dʒ, l-coda→w, rr→x, nasal om→õ"
                "   →   Regras produtivas do PT-BR genuinamente aprendidas",
                font_size=17, bold=True, color=BRANCO)

    footer(s, page_num=str(num))
    return s


def slide_generalization_overview(prs: Presentation, num: int, slide_data=None):
    """Slide 17 — Visão geral da generalização"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Resultados de Generalização — Visão Geral")

    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows = _table_from_parsed(slide_data["tables"][0])
    else:
        rows = [
            ["Categoria", "Corretas", "Score Fonol.", "Conclusão"],
            ["PT-BR reais OOV", "5/5", "100%", "Aprendeu regras, não memorizou"],
            ["Controles",         "3/4", "98%", "1 erro: ɣ→x (erro sistemático)"],
            ["Generalização PT-BR", "4/9", "97%", "Near-misses articulatórios"],
            ["Consoantes duplas",   "1/5", "81%", "Geminadas estrangeiras: gap"],
            ["Anglicismos",         "1/5", "71%", "Fonologia inglesa é OOV"],
            ["Chars OOV",           "0/3", "68%", "Falha esperada e documentada"],
        ]

    col_widths = [Inches(3.5), Inches(1.8), Inches(2.2), Inches(4.8)]
    add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), rows,
              col_widths=col_widths, font_size=15,
              highlight_rows={1})

    add_rect(s, Inches(0.5), Inches(5.55), W - Inches(1.0), Inches(0.7),
             fill_color=AZUL_CLARO)
    add_textbox(s, Inches(0.7), Inches(5.6), W - Inches(1.4), Inches(0.6),
                "Score fonológico: mesmo quando erra, o fonema está na família articulatória certa"
                " — 97% = near-miss, não falha total",
                font_size=15, color=AZUL_ESCURO)

    footer(s, page_num=str(num))
    return s


def slide_systematic_error(prs: Presentation, num: int):
    """Slide 18 — Erro sistemático ɣ→x"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Erro Sistemático Identificado",
                  "ɣ (velar vozeada) onde deveria ser x (velar surda) em coda")

    add_code_box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.2),
                 'borboleta   predição: b o ɣ b o ˈ l e t ə\n'
                 '            esperado: b o x b o ˈ l ɛ t ə   ← ɣ→x (vozeamento errado)\n\n'
                 'açucarzão   predição: a s u k a ɣ ˈ z ã ʊ̃\n'
                 '            esperado: a s u k a x ˈ z ã ʊ̃   ← ɣ→x (vozeamento errado)',
                 font_size=15)

    # Análise
    txb = s.shapes.add_textbox(Inches(0.5), Inches(3.9), W - Inches(1.0), Inches(2.5))
    tf = txb.text_frame
    tf.word_wrap = True
    add_para(tf, "O que isso revela:", font_size=17, bold=True, color=AZUL_ESCURO)
    add_para(tf, "• O modelo aprendeu a posição de coda → sabe que algo especial acontece ali",
             font_size=15, color=PRETO, space_before=6)
    add_para(tf, "• Mas não aprendeu a restrição fonotática: coda velar em PT-BR é tipicamente surda (/x/, não /ɣ/)",
             font_size=15, color=PRETO, space_before=4)
    add_para(tf, "• Métricas globais (PER/WER) não capturam isso — só análise qualitativa revela o padrão",
             font_size=15, color=CINZA_TEXTO, space_before=4)

    add_rect(s, Inches(0.5), Inches(6.35), W - Inches(1.0), Inches(0.7),
             fill_color=AZUL_CLARO)
    add_textbox(s, Inches(0.7), Inches(6.4), W - Inches(1.4), Inches(0.6),
                "Solução futura: integrar restrições fonotáticas PT-BR (Hayes & Wilson 2008)",
                font_size=15, color=AZUL_ESCURO)

    footer(s, page_num=str(num))
    return s


def slide_phonological_score(prs: Presentation, num: int):
    """Slide 19 — A Métrica Fonológica"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "A Métrica Fonológica (Score)",
                  "Distância articulatória independente do G2P — não usa PanPhon")

    add_code_box(s, Inches(0.5), Inches(1.45), Inches(6.4), Inches(1.5),
                 '_group_dist("ɣ", "x") = 0.1   # mesma família: FR_velar\n'
                 '_group_dist("a",  "k") = 0.9   # opostos: V_baixo vs OC_velar',
                 font_size=15)

    rows = [
        ["Score", "Rótulo", "Significado"],
        ["100%", "exato", "Transcrição idêntica"],
        ["90–99%", "muito próximo", "1 substituição na mesma família"],
        ["70–89%", "próximo", "Substituições em famílias relacionadas"],
        ["50–69%", "parcial", "Diferenças estruturais"],
        ["< 50%", "distante", "Falha substancial"],
    ]
    col_widths = [Inches(1.8), Inches(3.0), Inches(5.8)]
    add_table(s, Inches(0.5), Inches(3.1), Inches(10.6), rows,
              col_widths=col_widths, font_size=15,
              highlight_rows={1}, highlight_color=VERDE_CLARO,
              highlight_text_color=VERDE_SOTA)

    add_textbox(s, Inches(0.5), Inches(6.05), W - Inches(1.0), Inches(0.5),
                "Uso: G2PPredictor._phonological_score(pred, ref)  — independente de treino",
                font_size=14, color=CINZA_TEXTO, font_name="Courier New")

    footer(s, page_num=str(num))
    return s


def slide_usage(prs: Presentation, num: int, slide_data=None):
    """Slide 20 — Como usar"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Como Usar — 3 Comandos Essenciais")

    # Se slide_data fornecido e tem código, usa Markdown; senão hardcoded
    if slide_data and slide_data.get("code_blocks"):
        # Esperamos 3 blocos de código (1 por comando)
        code_blocks = slide_data["code_blocks"][:3]
        titles = ["1. Testar uma palavra", "2. Avaliar banco de generalização", "3. Modo interativo"]
        y = Inches(1.5)
        for title, code in zip(titles, code_blocks):
            add_textbox(s, Inches(0.5), y, W - Inches(1.0), Inches(0.4),
                        title, font_size=17, bold=True, color=AZUL_ESCURO)
            add_code_box(s, Inches(0.5), y + Inches(0.42), W - Inches(1.0),
                         Inches(0.95), code, font_size=15)
            y += Inches(1.7)
    else:
        cmds = [
            ("1. Testar uma palavra",
             "python src/inference_light.py --index 18 --word computador\n"
             "# → k õ p u t a ˈ d o x"),
            ("2. Avaliar banco de generalização",
             "python src/inference_light.py --index 18 --neologisms docs/generalization_test.tsv"),
            ("3. Modo interativo",
             "python src/inference_light.py --index 18 --interactive"),
        ]
        y = Inches(1.5)
        for title, code in cmds:
            add_textbox(s, Inches(0.5), y, W - Inches(1.0), Inches(0.4),
                        title, font_size=17, bold=True, color=AZUL_ESCURO)
            add_code_box(s, Inches(0.5), y + Inches(0.42), W - Inches(1.0),
                         Inches(0.7) if "\n" not in code else Inches(0.95),
                         code, font_size=15)
            y += Inches(1.7)

    add_rect(s, Inches(0.5), Inches(6.25), W - Inches(1.0), Inches(0.75),
             fill_color=AZUL_CLARO)
    add_textbox(s, Inches(0.7), Inches(6.3), W - Inches(1.4), Inches(0.65),
                "--index 18 = Exp104b (PER SOTA 0,49%)  ·  --index 11 = Exp9 (WER SOTA 4,96%)",
                font_size=15, color=AZUL_ESCURO, font_name="Calibri")

    footer(s, page_num=str(num))
    return s


def slide_model_choice(prs: Presentation, num: int):
    """Slide 21 — Escolha do modelo"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Escolha do Modelo na Prática")

    add_code_box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.5),
                 'Você quer...\n'
                 '    ├─ Acertar mais palavras inteiras?\n'
                 '    │   └─ Use Exp9 (index=11)   →   WER 4,96%\n'
                 '    │       ex: indexação de texto, busca fonética, NLP\n'
                 '    │\n'
                 '    └─ Acertar mais fonemas individuais?\n'
                 '        └─ Use Exp104b (index=18)   →   PER 0,49%\n'
                 '            ex: TTS, alinhamento fonético, análise linguística',
                 font_size=16)

    txb = s.shapes.add_textbox(Inches(0.5), Inches(5.2), W - Inches(1.0), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = True
    limits = [
        ("PT-BR novas", "Ambos generalizam bem"),
        ("Anglicismos / geminadas", "Espere erros — além do corpus"),
        ("k / w / y", "Mapeamento para <UNK> — avise o usuário"),
    ]
    for lim, desc in limits:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        r1 = p.add_run(); r1.text = f"  {lim}: "; r1.font.bold = True
        r1.font.size = Pt(15); r1.font.color.rgb = AZUL_ESCURO
        r2 = p.add_run(); r2.text = desc
        r2.font.size = Pt(15); r2.font.color.rgb = CINZA_TEXTO

    footer(s, page_num=str(num))
    return s


def slide_sota(prs: Presentation, num: int, slide_data=None):
    """Slide 22 — SOTA"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Comparação com Estado da Arte")

    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows = _table_from_parsed(slide_data["tables"][0])
    else:
        rows = [
            ["Sistema", "PER", "WER", "Idioma", "Teste (N)", "Params"],
            ["**FG2P Exp104b**", "**0,49%**", "5,43%", "PT-BR", "28.782", "9,7M"],
            ["**FG2P Exp9**",    "0,58%", "**4,96%**", "PT-BR", "28.782", "9,7M"],
            ["LatPhon 2025",     "0,86%", "—", "PT-BR", "500",    "n/d"],
            ["ByT5-Small",       "8,9%",  "—", "100 idiomas", "~500/lang", "299M"],
        ]

    col_widths = [Inches(3.2), Inches(1.5), Inches(1.5), Inches(2.0), Inches(2.3), Inches(1.8)]
    add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), rows,
              col_widths=col_widths, font_size=15,
              highlight_rows={1, 2})

    add_rect(s, Inches(0.5), Inches(4.5), W - Inches(1.0), Inches(0.75),
             fill_color=AZUL_CLARO)
    add_textbox(s, Inches(0.7), Inches(4.55), W - Inches(1.4), Inches(0.65),
                "FG2P usa 9,7M params  ·  ByT5-Small usa 299M (30× maior, zero-shot)"
                "  ·  FG2P treinado especificamente em PT-BR",
                font_size=15, color=AZUL_ESCURO)

    txb = s.shapes.add_textbox(Inches(0.5), Inches(5.45), W - Inches(1.0), Inches(1.5))
    tf = txb.text_frame; tf.word_wrap = True
    add_para(tf, "PER: Phoneme Error Rate — % de fonemas errados (menor = melhor)",
             font_size=14, color=CINZA_TEXTO)
    add_para(tf, "WER: Word Error Rate — % de palavras com qualquer erro (menor = melhor)",
             font_size=14, color=CINZA_TEXTO, space_before=4)

    footer(s, page_num=str(num))
    return s


def slide_limits(prs: Presentation, num: int, slide_data=None):
    """Slide 23 — Limites"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Limites Bem-Definidos")

    # Se slide_data fornecido e tem tabelas, usa dados do Markdown
    if slide_data and slide_data.get("tables"):
        rows = _table_from_parsed(slide_data["tables"][0])
    else:
        rows = [
            ["Limite", "Causa", "Solução Futura"],
            ["Homógrafos heterófonos (jogo, gosto)", "Ambiguidade léxico-sintática sem contexto", "Pipeline NLP→G2P em série"],
            ["Geminadas (zz, pp, tt)", "Sem exemplos no corpus", "Ampliar corpus com empréstimos"],
            ["Fonologia inglesa", "OOV fonológico", "Corpus bilíngue adaptado"],
            ["k, w, y", "OOV de charset", "Re-treinar com charset ampliado"],
            ["ɣ→x em coda", "Sem restrição fonotática", "Integrar fonotática PT-BR"],
            [".↔ˈ confusão", "Posicional, não só distância", "Métricas separadas + beam restrito"],
        ]

    col_widths = [Inches(3.8), Inches(3.8), Inches(4.7)]
    add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), rows,
              col_widths=col_widths, font_size=15)

    add_rect(s, Inches(0.5), Inches(5.6), W - Inches(1.0), Inches(0.7),
             fill_color=AZUL_CLARO)
    add_textbox(s, Inches(0.7), Inches(5.65), W - Inches(1.4), Inches(0.6),
                "Limites documentados são uma virtude — permitem aplicação informada e desenvolvimento direcionado",
                font_size=15, color=AZUL_ESCURO)

    footer(s, page_num=str(num))
    return s


def slide_next_steps(prs: Presentation, num: int):
    """Slide 24 — Próximos passos"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Próximos Passos")

    col_specs = [
        ("Curto Prazo", VERDE_CLARO, VERDE_SOTA, [
            "Stress Accuracy: % de acentos na sílaba certa",
            "Boundary F1: precisão/revocação das fronteiras",
            "Corpus ampliado com empréstimos e geminadas",
        ]),
        ("Médio Prazo", AZUL_CLARO, AZUL_ESCURO, [
            "Espaço articulatório 7D contínuo",
            "Símbolos . e ˈ distinguíveis intrinsecamente",
            "Interpolação entre sons (morfologia fonética)",
        ]),
        ("Longo Prazo", RGBColor(0xf3, 0xe8, 0xff), RGBColor(0x7e, 0x22, 0xce), [
            "Universalização multilíngue via 7D",
            "Transfer: PT-BR → Espanhol → Inglês",
            "Fine-tuning com corpus reduzido",
        ]),
    ]
    col_w = Inches(4.1)
    for i, (title, bg, fg, items) in enumerate(col_specs):
        lft = Inches(0.4) + i * Inches(4.3)
        add_rect(s, lft, Inches(1.45), col_w, Inches(5.6), fill_color=bg)
        add_textbox(s, lft + Inches(0.15), Inches(1.5), col_w - Inches(0.3), Inches(0.5),
                    title, font_size=17, bold=True, color=fg)
        txb = s.shapes.add_textbox(lft + Inches(0.15), Inches(2.1),
                                   col_w - Inches(0.3), Inches(4.7))
        tf = txb.text_frame; tf.word_wrap = True
        for item in items:
            add_para(tf, f"• {item}", font_size=14, color=PRETO, space_before=8)

    footer(s, page_num=str(num))
    return s


def slide_summary(prs: Presentation, num: int, slide_data=None):
    """Slide 25 — Resumo em 5 pontos"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=BRANCO)
    add_title_bar(s, "Resumo em 5 Pontos")

    # Se slide_data tem prose, extrai 5 pontos; senão hardcoded
    if slide_data and slide_data.get("prose"):
        # Esperamos texto no formato "1. Título — Descrição"
        prose_text = slide_data["prose"]
        points = []
        for line in prose_text.split('\n'):
            line = line.strip()
            if line and line[0].isdigit() and '.' in line:
                # Formato: "1. Título — descrição" ou "1. **Título** — descrição"
                num_str = line[0]
                rest = line[3:].strip()  # Remove "1. "
                if "—" in rest:
                    title, desc = rest.split("—", 1)
                    title = title.strip().replace("**", "")
                    desc = desc.strip()
                    points.append((num_str, title, desc))
                elif "·" in rest:
                    # Alternativa: "1. Título · Descrição"
                    title, desc = rest.split("·", 1)
                    title = title.strip().replace("**", "")
                    desc = desc.strip()
                    points.append((num_str, title, desc))

        if not points:  # Fallback se formato inválido
            points = [
                ("1", "BiLSTM + Atenção",
                 "Arquitetura clássica, bem compreendida — SOTA em PT-BR com apenas 9,7M params"),
                ("2", "Distance-Aware Loss",
                 "Sinal fonológico λ=0,20: ensina o modelo a preferir erros inteligentes"),
                ("3", "Trade-off Pareto",
                 "Separadores: PER↓ + WER↑ de forma irredutível — escolha depende da aplicação"),
                ("4", "Custom Distances",
                 "Fix pós-normalização para . e ˈ (vetores zero) → SOTA PER 0,49%"),
                ("5", "Generalização Real",
                 "100% em palavras PT-BR OOV — falhas têm limites bem-definidos"),
            ]
    else:
        points = [
            ("1", "BiLSTM + Atenção",
             "Arquitetura clássica, bem compreendida — SOTA em PT-BR com apenas 9,7M params"),
            ("2", "Distance-Aware Loss",
             "Sinal fonológico λ=0,20: ensina o modelo a preferir erros inteligentes"),
            ("3", "Trade-off Pareto",
             "Separadores: PER↓ + WER↑ de forma irredutível — escolha depende da aplicação"),
            ("4", "Custom Distances",
             "Fix pós-normalização para . e ˈ (vetores zero) → SOTA PER 0,49%"),
            ("5", "Generalização Real",
             "100% em palavras PT-BR OOV — falhas têm limites bem-definidos"),
        ]

    for i, (num_str, title, desc) in enumerate(points):
        y = Inches(1.45) + i * Inches(1.05)
        add_rect(s, Inches(0.4), y, Inches(0.7), Inches(0.85),
                 fill_color=AZUL_ESCURO)
        add_textbox(s, Inches(0.4), y + Inches(0.15), Inches(0.7), Inches(0.6),
                    num_str, font_size=24, bold=True, color=BRANCO,
                    align=PP_ALIGN.CENTER)
        add_rect(s, Inches(1.2), y, Inches(11.5), Inches(0.85),
                 fill_color=AZUL_CLARO if i % 2 == 0 else CINZA_CLARO)
        add_textbox(s, Inches(1.35), y + Inches(0.05), Inches(3.0), Inches(0.75),
                    title, font_size=16, bold=True, color=AZUL_ESCURO)
        add_textbox(s, Inches(4.5), y + Inches(0.1), Inches(8.0), Inches(0.7),
                    desc, font_size=15, color=CINZA_TEXTO)

    footer(s, page_num=str(num))
    return s


def slide_end(prs: Presentation, num: int):
    """Slide 26 — Obrigado"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_color=AZUL_ESCURO)
    add_rect(s, 0, H - Inches(2.0), W, Inches(2.0),
             fill_color=RGBColor(0x17, 0x30, 0x78))

    add_textbox(s, Inches(0.5), Inches(1.5), W - Inches(1.0), Inches(1.5),
                "Obrigado", font_size=60, bold=True, color=BRANCO,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(3.0), W - Inches(1.0), Inches(0.6),
                "FG2P · PT-BR G2P BiLSTM Encoder-Decoder",
                font_size=22, color=RGBColor(0xba, 0xd3, 0xf8),
                align=PP_ALIGN.CENTER)

    add_code_box(s, Inches(2.5), Inches(3.9), Inches(8.3), Inches(0.9),
                 '"biscoitinhozão" → b i s k o y t ʃ ĩ ɲ ɔ ˈ z ã ʊ̃',
                 font_size=18)

    txb = s.shapes.add_textbox(Inches(0.5), Inches(5.1), W - Inches(1.0), Inches(1.0))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Código: src/  ·  Dados: dicts/  ·  Artigo: docs/15_ARTIGO.md"
    r.font.size = Pt(16); r.font.color.rgb = RGBColor(0x93, 0xc5, 0xfd)

    add_code_box(s, Inches(2.5), Inches(6.2), Inches(8.3), Inches(0.7),
                 'python src/inference_light.py --index 18 --interactive',
                 font_size=16)
    return s


def slide_glossario_fonetica(prs: Presentation, num: int, slide_data=None):
    """Slide 27 — APENDICE A: Articulacoes Vocálicas"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(0.8), fill_color=AZUL_ESCURO)

    add_textbox(s, Inches(0.5), Inches(0.15), W - Inches(1.0), Inches(0.5),
                "APENDICE A: Como os Sons Nascem na Boca",
                font_size=32, bold=True, color=BRANCO)

    # Conteudo em bullets (simplificado para caber no slide)
    content = [
        "Ponto de Articulacao: Labial (/p/, /b/, /m/), Alveolar (/t/, /d/, /s/), Palatal (/ʃ/, /ʒ/), Velar (/k/, /ɡ/)",
        "Modo: Oclusiva (bloqueia), Fricativa (atrito), Nasal (ar na nariz), Lateral (ar nos lados)",
        "Vozeamento: Vozeadas (/b/, /d/, /ɡ/) vibram; Desvozeadas (/p/, /t/, /k/) não vibram",
        "Vogais: Alto (/i/, /u/), Medio (/e/, /o/, /ə/), Baixo (/a/)",
    ]

    top = Inches(1.2)
    for bullet in content:
        add_textbox(s, Inches(0.7), top, W - Inches(1.4), Inches(0.6),
                    f"• {bullet}", font_size=14, color=CINZA_TEXTO, wrap=True)
        top += Inches(0.65)

    return s


def slide_glossario_algoritmos(prs: Presentation, num: int, slide_data=None):
    """Slide 28 — APENDICE B: Termos de Algoritmos"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(0.8), fill_color=AZUL_ESCURO)

    add_textbox(s, Inches(0.5), Inches(0.15), W - Inches(1.0), Inches(0.5),
                "APENDICE B: Termos de Algoritmos e ML",
                font_size=32, bold=True, color=BRANCO)

    # ML Basics em tabela simples (header + 4 rows)
    rows = [
        ["Termo", "Definicao"],
        ["Modelo", "Funcao que aprende padroes dos dados"],
        ["Treino", "Ajustar parametros para minimizar erro"],
        ["Validacao", "Dados para monitorar progresso"],
        ["Teste", "Dados nunca vistos para avaliar"],
    ]
    add_table(s, Inches(0.7), Inches(1.2), Inches(11.9), rows,
              header_fill=AZUL_MEDIO, row_alt_fill=CINZA_CLARO)

    # Metricas (abaixo)
    add_textbox(s, Inches(0.7), Inches(3.9), W - Inches(1.4), Inches(0.35),
                "Metricas: PER (Phoneme Error Rate) - WER (Word Error Rate) - Accuracy",
                font_size=12, color=CINZA_TEXTO)

    add_textbox(s, Inches(0.7), Inches(4.3), W - Inches(1.4), Inches(0.35),
                "Arquitetura: LSTM (memoria longa) - Attention (foca em partes) - Embedding (converte simbolos)",
                font_size=12, color=CINZA_TEXTO)

    add_textbox(s, Inches(0.7), Inches(4.7), W - Inches(1.4), Inches(0.35),
                "Losses: Cross Entropy (CE) - Distance-Aware (DA) com penalidade proporcional a distancia",
                font_size=12, color=CINZA_TEXTO)

    return s


def slide_glossario_projeto(prs: Presentation, num: int, slide_data=None):
    """Slide 29 — APENDICE C: Termos do Projeto"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(0.8), fill_color=AZUL_ESCURO)

    add_textbox(s, Inches(0.5), Inches(0.15), W - Inches(1.0), Inches(0.5),
                "APENDICE C: Termos do Projeto e Fenomenos Linguisticos",
                font_size=32, bold=True, color=BRANCO)

    # Dois blocos: Conceitos e Fenomenos
    add_textbox(s, Inches(0.7), Inches(1.1), W - Inches(1.4), Inches(0.35),
                "Conceitos Principais",
                font_size=14, bold=True, color=AZUL_MEDIO)

    content_left = [
        "G2P: Converter grafemas (letras) em fonemas (sons)",
        "PT-BR: Portugues brasileiro — variante que modelamos",
        "SOTA: State-of-the-Art — melhor resultado ate agora",
    ]

    top = Inches(1.55)
    for bullet in content_left:
        add_textbox(s, Inches(0.9), top, W - Inches(2.0), Inches(0.45),
                    f"• {bullet}", font_size=12, color=CINZA_TEXTO, wrap=True)
        top += Inches(0.5)

    add_textbox(s, Inches(0.7), Inches(3.2), W - Inches(1.4), Inches(0.35),
                "Fenomenos Linguisticos",
                font_size=14, bold=True, color=AZUL_MEDIO)

    content_right = [
        "Coda: Posicao final de silaba — consoantes sofrem mudancas",
        "Stress: Silaba tonica vs. atona — acentuacao",
        "Reducao Vocálica: Vogal atona muda de som",
        "Palatalizacao: /t/, /d/ → /tʃ/, /dʒ/ antes de /i/",
    ]

    top = Inches(3.65)
    for bullet in content_right:
        add_textbox(s, Inches(0.9), top, W - Inches(2.0), Inches(0.45),
                    f"• {bullet}", font_size=12, color=CINZA_TEXTO, wrap=True)
        top += Inches(0.5)

    return s


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_presentation_compact_from_markdown(parsed: dict, output_path: Path, meta: dict) -> Path:
    """Gera PPTX COMPACTO (10 minutos / ~13 slides) usando funções hardcoded selecionadas.

    Versão reduzida que mantém qualidade visual usando as mesmas funções da versão completa,
    apenas selecionando os slides essenciais para 10 minutos de apresentação.
    """
    slides_data = parsed.get("slides", [])
    prs = new_prs()

    logger.info("Gerando slides (modo Compacto)...")

    # Slides essenciais para versão compacta (10 minutos)
    slide_opening(prs, slide_data=meta)                              # 0: Abertura
    slide_title(prs)                                                  # 1: Título
    slide_g2p_problem(prs, 2)                                        # 2: O que é G2P
    slide_ptbr_hard(prs, 3, slide_data=slides_data[1] if len(slides_data) > 1 else None)  # 3: Por que difícil
    slide_data(prs, 4, slide_data=slides_data[2] if len(slides_data) > 2 else None)       # 4: Dados + Arquitetura
    slide_da_loss_problem(prs, 5)                                    # 5: Distance-Aware Loss - Problema (hardcoded)
    slide_da_loss(prs, 6, slide_data=slides_data[4] if len(slides_data) > 4 else None)   # 6: Distance-Aware Loss - Fórmula
    slide_evolution(prs, 7)                                          # 7: Experimentos (hardcoded)
    slide_oov_result(prs, 8)                                         # 8: Generalização (hardcoded only)
    slide_systematic_error(prs, 9)                                   # 9: Descoberta Fonológica (hardcoded)
    slide_usage(prs, 10)                                             # 10: Como usar (hardcoded)
    slide_sota(prs, 11)                                              # 11: SOTA + Limites (hardcoded)
    slide_summary(prs, 12)                                           # 12: Resumo (hardcoded)
    slide_end(prs, 13)                                               # 13: Obrigado

    # Apêndices de glossário (mantidos da versão completa)
    slide_glossario_fonetica(prs, 14)                                # 14: Glossário Fonética
    slide_glossario_algoritmos(prs, 15)                              # 15: Glossário Algoritmos
    slide_glossario_projeto(prs, 16)                                 # 16: Glossário Termos

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info(f"Apresentação compacta (Markdown) salva em: {output_path}")
    return output_path


def build_presentation_from_markdown(md_path: Path, output_path: Path, filtered_markdown: str = None) -> Path:
    """Gera PPTX usando conteúdo do Markdown com suporte a modo unificado.

    Suporta:
    - Versão COMPLETA (26 slides + 3 apêndices = 29 total)
    - Versão COMPACTA (10-11 slides + 3 apêndices = 13-14 total)
    - Modo UNIFICADO: [modes: ...] tags automáti camente filtradas

    Args:
        md_path: Caminho para arquivo Markdown
        output_path: Caminho de saída .pptx
        filtered_markdown: (opcional) Conteúdo MD já filtrado pelo modo
    """
    if parse_presentation is None:
        raise ImportError("presentation_parser não encontrado. Verifique src/reporting/presentation_parser.py")

    logger.info(f"Lendo apresentação de: {md_path}")

    # Se markdown filtrado foi passado, usar; senão ler arquivo
    if filtered_markdown is not None:
        # Salvar em arquivo temporário para parse_presentation
        import tempfile
        with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as tmp:
            tmp.write(filtered_markdown)
            tmp_path = Path(tmp.name)
        try:
            parsed = parse_presentation(tmp_path)
        finally:
            tmp_path.unlink()  # Limpar arquivo temporário
    else:
        parsed = parse_presentation(md_path)
    meta = parsed.get("meta", {})
    slides_data = parsed.get("slides", [])

    # DETECTAR VERSÃO / MODO (compatível com arquivo MERGED)
    # Em modo unified, o filtered_markdown já foi processado, logo sabemos o mode pelo tamanho
    is_compact = (meta.get("presentation_type") == "compact" or
                  "COMPACTA" in str(md_path).upper() or
                  len(slides_data) <= 20)  # Heurística: compact tem ~19 slides

    if is_compact:
        logger.info(f"  [MODO COMPACTO] Versão simplificada (10 minutos)")
        return build_presentation_compact_from_markdown(parsed, output_path, meta)
    logger.info(f"  author: {meta.get('author', '?')}")
    logger.info(f"  professor: {meta.get('professor', '?')}")
    logger.info(f"  {len(slides_data)} slides encontrados")

    prs = new_prs()
    logger.info("Gerando slides (modo Markdown)...")
    slide_opening(prs, slide_data=meta)
    slide_title(prs)
    slide_g2p_problem(prs, 2)
    slide_ptbr_hard(prs, 3, slide_data=slides_data[2] if len(slides_data) > 2 else None)
    slide_data(prs, 4, slide_data=slides_data[3] if len(slides_data) > 3 else None)
    slide_architecture(prs, 5)
    slide_attention(prs, 6)
    slide_ce_problem(prs, 7)
    slide_da_loss(prs, 8, slide_data=slides_data[7] if len(slides_data) > 7 else None)
    slide_da_example(prs, 9)
    slide_confidence_factor(prs, 10)
    slide_separators(prs, 11, slide_data=slides_data[10] if len(slides_data) > 10 else None)
    slide_custom_dist(prs, 12, slide_data=slides_data[11] if len(slides_data) > 11 else None)
    slide_evolution(prs, 13)
    slide_ranking(prs, 14, slide_data=slides_data[13] if len(slides_data) > 13 else None)
    slide_generalization_design(prs, 15, slide_data=slides_data[15] if len(slides_data) > 15 else None)
    slide_oov_result(prs, 16, slide_data=slides_data[16] if len(slides_data) > 16 else None)
    slide_generalization_overview(prs, 17, slide_data=slides_data[17] if len(slides_data) > 17 else None)
    slide_systematic_error(prs, 18)
    slide_phonological_score(prs, 19)
    slide_usage(prs, 20, slide_data=slides_data[19] if len(slides_data) > 19 else None)
    slide_model_choice(prs, 21)
    slide_sota(prs, 22, slide_data=slides_data[21] if len(slides_data) > 21 else None)
    slide_limits(prs, 23, slide_data=slides_data[22] if len(slides_data) > 22 else None)
    slide_next_steps(prs, 24)
    slide_summary(prs, 25, slide_data=slides_data[24] if len(slides_data) > 24 else None)
    slide_end(prs, 26)

    # Glossario slides (27-29)
    slide_glossario_fonetica(prs, 27, slide_data=slides_data[25] if len(slides_data) > 25 else None)
    slide_glossario_algoritmos(prs, 28, slide_data=slides_data[26] if len(slides_data) > 26 else None)
    slide_glossario_projeto(prs, 29, slide_data=slides_data[27] if len(slides_data) > 27 else None)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info(f"Apresentação (Markdown) salva em: {output_path}")
    return output_path


def build_presentation(output_path: Path) -> Path:
    prs = new_prs()

    logger.info("Gerando slides...")
    slide_opening(prs)  # Página de abertura personalizada (fallback hardcoded)
    slide_title(prs)
    slide_g2p_problem(prs, 2)
    slide_ptbr_hard(prs, 3)
    slide_data(prs, 4)
    slide_architecture(prs, 5)
    slide_attention(prs, 6)
    slide_ce_problem(prs, 7)
    slide_da_loss(prs, 8)
    slide_da_example(prs, 9)
    slide_confidence_factor(prs, 10)
    slide_separators(prs, 11)
    slide_custom_dist(prs, 12)
    slide_evolution(prs, 13)
    slide_ranking(prs, 14)
    slide_generalization_design(prs, 15)
    slide_oov_result(prs, 16)
    slide_generalization_overview(prs, 17)
    slide_systematic_error(prs, 18)
    slide_phonological_score(prs, 19)
    slide_usage(prs, 20)
    slide_model_choice(prs, 21)
    slide_sota(prs, 22)
    slide_limits(prs, 23)
    slide_next_steps(prs, 24)
    slide_summary(prs, 25)
    slide_end(prs, 26)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info(f"Apresentação salva em: {output_path}")
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Gera apresentação PPTX para FG2P (Unified Mode System)",
        epilog="Exemplos:\n"
               "  python presentation_generator.py --mode full    # 29 slides\n"
               "  python presentation_generator.py --mode compact # 16 slides (10 min)\n"
               "  python presentation_generator.py --mode full -o custom_output.pptx",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "--mode",
        choices=["full", "compact"],
        default="full",
        help="Modo de apresentação: 'full' (29 slides) ou 'compact' (16 slides, 10 min). Padrão: full",
    )
    parser.add_argument(
        "--output", "-o",
        type=Path,
        default=RESULTS_DIR / "fg2p_presentation.pptx",
        help="Caminho de saída para o arquivo .pptx (padrão: results/fg2p_presentation.pptx)",
    )
    parser.add_argument(
        "--markdown", "-m",
        type=Path,
        metavar="MD_PATH",
        help="Arquivo Markdown customizado (padrão: docs/17_APRESENTACAO.md)",
    )
    # Backward compatibility flags (deprecated)
    parser.add_argument(
        "--compact",
        action="store_true",
        help="[DEPRECATED] Use --mode compact em vez disso",
    )
    parser.add_argument(
        "--duration",
        type=int,
        metavar="MINUTES",
        help="[DEPRECATED] Duração em minutos. Use --mode em vez disso.",
    )
    args = parser.parse_args()

    # Backward compatibility: converter flags antigos para novo sistema
    mode = args.mode
    if args.compact:
        logger.warning("[DEPRECATED] Use --mode compact em vez de --compact")
        mode = "compact"
    if args.duration:
        logger.warning("[DEPRECATED] Use --mode em vez de --duration")
        if args.duration <= 7:
            mode = "compact"
        elif args.duration <= 12:
            mode = "compact"
        else:
            mode = "full"

    # Selecionar arquivo Markdown
    if args.markdown:
        markdown_file = args.markdown
    else:
        # Usar arquivo canonical com modo filtering
        markdown_file = Path(__file__).parent.parent.parent / "docs" / "17_APRESENTACAO.md"

    if not markdown_file.exists():
        logger.error(f"Arquivo Markdown não encontrado: {markdown_file}")
        sys.exit(1)

    # Carregar e filtrar Markdown pelo mode
    logger.info(f"[MODE] {mode.upper()} | {markdown_file.name}")
    with open(markdown_file, encoding='utf-8') as f:
        markdown_content = f.read()

    # Aplicar filtro de modo
    filtered_markdown = filter_markdown_by_mode(markdown_content, mode)

    # ÚNICO modo: Markdown-driven (sem redundância hardcoded)
    out = build_presentation_from_markdown(markdown_file, args.output, filtered_markdown)
    print(f"\n[OK] Apresentacao gerada: {out}")
    print(f"   Mode: {mode.upper()} | Arquivo: {markdown_file.name}")

    # Calcular número de slides (do arquivo filtrado)
    slide_count = filtered_markdown.count('\n---\n') + 1
    duration_str = "10 minutos (compacto)" if mode == "compact" else "29 slides (completo)"

    print(f"   {slide_count} slides · {duration_str} · Tema azul académico · 16:9 widescreen")
    print(f"   Fonte: {markdown_file.name}")
    print(f"\nAbrir com: Microsoft PowerPoint / LibreOffice Impress / Google Slides")


if __name__ == "__main__":
    main()
